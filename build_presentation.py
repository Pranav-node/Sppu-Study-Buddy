import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# Load original template
template_path = r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx'
prs = pptx.Presentation(template_path)

# Helper to clone slide with shapes
def clone_slide(prs, source_slide):
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for sp in list(new_slide.shapes):
        sp_elem = sp.element
        sp_elem.getparent().remove(sp_elem)
    for sp in source_slide.shapes:
        new_sp = parse_xml(sp.element.xml)
        new_slide.shapes._spTree.append(new_sp)
    return new_slide

# Colors matching template
COLOR_RED = RGBColor(192, 0, 0)        # #C00000
COLOR_NAVY = RGBColor(0, 32, 96)       # #002060
COLOR_MAGENTA = RGBColor(204, 0, 102)  # #CC0066
COLOR_PURPLE = RGBColor(102, 0, 204)   # #6600CC
COLOR_DARK_RED = RGBColor(153, 0, 0)   # #990000
COLOR_BODY = RGBColor(40, 40, 40)      # #282828
COLOR_CARD_BG = RGBColor(248, 249, 250)
COLOR_CARD_BORDER = RGBColor(210, 215, 220)
COLOR_ORANGE_BG = RGBColor(254, 243, 235)
COLOR_ACCENT_BLUE = RGBColor(24, 76, 120)

# Slide 1: Update existing Slide 1
slide1 = prs.slides[0]
for sp in slide1.shapes:
    if sp.is_placeholder and sp.placeholder_format.type == pptx.enum.shapes.PP_PLACEHOLDER.SUBTITLE:
        tf = sp.text_frame
        tf.clear()
        
        # P0: Title
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0_lbl = p0.add_run()
        r0_lbl.text = "Project Title: "
        r0_lbl.font.name = "Georgia"
        r0_lbl.font.size = Pt(20)
        r0_lbl.font.color.rgb = COLOR_RED
        
        r0_txt = p0.add_run()
        r0_txt.text = "AI Voice Support Agent Using Speech Recognition and Large Language Models"
        r0_txt.font.name = "Georgia"
        r0_txt.font.size = Pt(20)
        r0_txt.font.bold = True
        r0_txt.font.color.rgb = COLOR_RED

        # Empty lines
        for _ in range(2):
            p = tf.add_paragraph()
            p.text = ""

        # P: Group No.
        p_grp = tf.add_paragraph()
        p_grp.alignment = PP_ALIGN.CENTER
        r_grp = p_grp.add_run()
        r_grp.text = "Group No.: [Your Group Number]"
        r_grp.font.name = "Georgia"
        r_grp.font.size = Pt(18)
        r_grp.font.bold = True
        r_grp.font.color.rgb = COLOR_NAVY

        # Empty line
        p = tf.add_paragraph()
        p.text = ""

        # P: Group Members
        p_mem = tf.add_paragraph()
        p_mem.alignment = PP_ALIGN.CENTER
        r_mem = p_mem.add_run()
        r_mem.text = "Group Members: [Add Member Names with Roll Numbers]"
        r_mem.font.name = "Georgia"
        r_mem.font.size = Pt(18)
        r_mem.font.bold = True
        r_mem.font.color.rgb = COLOR_PURPLE

        # Empty lines
        for _ in range(2):
            p = tf.add_paragraph()
            p.text = ""

        # P: Project Guide
        p_gd1 = tf.add_paragraph()
        p_gd1.alignment = PP_ALIGN.CENTER
        r_gd1 = p_gd1.add_run()
        r_gd1.text = "Project Guide"
        r_gd1.font.name = "Georgia"
        r_gd1.font.size = Pt(18)
        r_gd1.font.color.rgb = COLOR_DARK_RED

        p_gd2 = tf.add_paragraph()
        p_gd2.alignment = PP_ALIGN.CENTER
        r_gd2 = p_gd2.add_run()
        r_gd2.text = "[Guide Name]"
        r_gd2.font.name = "Georgia"
        r_gd2.font.size = Pt(18)
        r_gd2.font.bold = True
        r_gd2.font.color.rgb = COLOR_DARK_RED

# Slide 2: Contents
slide2 = prs.slides[1]
# Text is already updated/verified, but let's re-set clean formatting
for sp in slide2.shapes:
    if sp.name.startswith("Google Shape;63"): # Title
        tf = sp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "Contents"
        r.font.name = "Times New Roman"
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = COLOR_RED
    elif sp.name.startswith("Google Shape;64"): # Body
        tf = sp.text_frame
        tf.clear()
        contents_items = [
            "1. Problem Statement",
            "2. Motivation",
            "3. Objectives",
            "4. Literature Survey",
            "5. Project Overview"
        ]
        for idx, item in enumerate(contents_items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = item
            r.font.name = "Times New Roman"
            r.font.size = Pt(24)
            r.font.bold = True
            r.font.color.rgb = COLOR_MAGENTA

# Store slide2 as template for content slides
slide_template = slide2

# Remove original slide 3 (Thank you) for now, we will add content slides and put Thank You at the end!
# Actually, let's keep slide3 oxml or clone from slide_template for all remaining slides!
rId = prs.slides._sldIdLst[2].rId
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[2]

# Function to add standard slide header/title
def create_base_slide(prs, title_text):
    s = clone_slide(prs, slide_template)
    # Clear body text box shape
    body_sp = None
    for sp in s.shapes:
        if sp.name.startswith("Google Shape;63"): # Title placeholder
            tf = sp.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = title_text
            r.font.name = "Times New Roman"
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = COLOR_RED
        elif sp.name.startswith("Google Shape;64"): # Body text box
            body_sp = sp
    if body_sp:
        # remove default body text box element so we can add custom structured boxes
        sp_elem = body_sp.element
        sp_elem.getparent().remove(sp_elem)
    return s

# ----------------------------------------------------
# SLIDE 3: Problem Statement
# ----------------------------------------------------
slide3 = create_base_slide(prs, "Problem Statement")

# Card 1: Problem (Left side)
left_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(4.2), Inches(3.4))
left_card.fill.solid()
left_card.fill.fore_color.rgb = RGBColor(253, 242, 242) # soft red tint
left_card.line.color.rgb = COLOR_RED
left_card.line.width = Pt(1.5)

tf = left_card.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
tf.word_wrap = True

p_hdr = tf.paragraphs[0]
p_hdr.alignment = PP_ALIGN.LEFT
r = p_hdr.add_run()
r.text = "⚠️ Problem"
r.font.name = "Georgia"
r.font.size = Pt(20)
r.font.bold = True
r.font.color.rgb = COLOR_RED

prob_bullets = [
    "Traditional customer support relies on rigid IVR menus or human agents.",
    "Results in long waiting times, high operational costs, and inconsistent experiences.",
    "Most existing systems fail to understand natural spoken language and real-time intent."
]
for b in prob_bullets:
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = "• " + b
    r.font.name = "Times New Roman"
    r.font.size = Pt(15)
    r.font.color.rgb = COLOR_BODY

# Card 2: Proposed Solution (Right side)
right_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.2), Inches(3.4))
right_card.fill.solid()
right_card.fill.fore_color.rgb = RGBColor(240, 248, 255) # soft blue tint
right_card.line.color.rgb = COLOR_NAVY
right_card.line.width = Pt(1.5)

tf = right_card.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
tf.word_wrap = True

p_hdr = tf.paragraphs[0]
p_hdr.alignment = PP_ALIGN.LEFT
r = p_hdr.add_run()
r.text = "💡 Proposed Solution"
r.font.name = "Georgia"
r.font.size = Pt(20)
r.font.bold = True
r.font.color.rgb = COLOR_NAVY

sol_bullets = [
    "Develop an AI-powered Voice Support Agent enabling real-time continuous voice interaction.",
    "Converts customer speech to text accurately using advanced STT models.",
    "Understands intent and generates responses using a Large Language Model (LLM).",
    "Synthesizes natural speech output for a human-like conversational experience."
]
for b in sol_bullets:
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(6)
    r = p.add_run()
    r.text = "• " + b
    r.font.name = "Times New Roman"
    r.font.size = Pt(14)
    r.font.color.rgb = COLOR_BODY

# Add icon decoration top left of cards
slide3.shapes.add_picture(r'd:\B.E final year project\icons\mic_icon.png', Inches(4.2), Inches(1.9), Inches(0.4), Inches(0.4))
slide3.shapes.add_picture(r'd:\B.E final year project\icons\chatbot_icon.png', Inches(8.8), Inches(1.9), Inches(0.4), Inches(0.4))


# ----------------------------------------------------
# SLIDE 4: Motivation
# ----------------------------------------------------
slide4 = create_base_slide(prs, "Motivation")

# Main content container card
mot_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(8.8), Inches(3.4))
mot_card.fill.solid()
mot_card.fill.fore_color.rgb = COLOR_CARD_BG
mot_card.line.color.rgb = COLOR_CARD_BORDER
mot_card.line.width = Pt(1)

tf = mot_card.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.25)
tf.word_wrap = True

mot_points = [
    ("Critical Business Requirement:", " Customer service is essential for every modern enterprise, but human-agent support is expensive and difficult to scale."),
    ("AI & NLP Breakthroughs:", " Recent advancements in Artificial Intelligence, Speech Recognition, and Large Language Models enable human-like voice agents."),
    ("Instant & 24/7 Availability:", " Eliminates long customer waiting queues by providing immediate, accurate, round-the-clock voice assistance."),
    ("Cost Efficiency & UX:", " Our motivation is to develop a smart AI Voice Support Agent that drastically improves customer satisfaction while lowering operational costs.")
]

for idx, (title, desc) in enumerate(mot_points):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(10)
    
    r_t = p.add_run()
    r_t.text = "✦ " + title
    r_t.font.name = "Times New Roman"
    r_t.font.size = Pt(17)
    r_t.font.bold = True
    r_t.font.color.rgb = COLOR_NAVY
    
    r_d = p.add_run()
    r_d.text = desc
    r_d.font.name = "Times New Roman"
    r_d.font.size = Pt(16)
    r_d.font.color.rgb = COLOR_BODY

slide4.shapes.add_picture(r'd:\B.E final year project\icons\ai_icon.png', Inches(8.8), Inches(1.9), Inches(0.45), Inches(0.45))


# ----------------------------------------------------
# SLIDE 5: Objectives
# ----------------------------------------------------
slide5 = create_base_slide(prs, "Objectives")

obj_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(8.8), Inches(3.4))
obj_card.fill.solid()
obj_card.fill.fore_color.rgb = COLOR_CARD_BG
obj_card.line.color.rgb = COLOR_CARD_BORDER
obj_card.line.width = Pt(1)

tf = obj_card.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
tf.word_wrap = True

objectives = [
    "Develop a real-time voice-based customer support system.",
    "Convert user speech into text accurately using speech recognition.",
    "Generate intelligent, context-aware responses using an AI language model.",
    "Convert generated text into natural, human-like speech output.",
    "Reduce customer waiting time and eliminate support queue bottlenecks.",
    "Improve service availability through 24×7 automated assistance.",
    "Provide a highly scalable and cost-effective customer support solution."
]

for idx, obj in enumerate(objectives):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(5)
    
    r_num = p.add_run()
    r_num.text = f"  {idx+1}. "
    r_num.font.name = "Times New Roman"
    r_num.font.size = Pt(16)
    r_num.font.bold = True
    r_num.font.color.rgb = COLOR_MAGENTA
    
    r_txt = p.add_run()
    r_txt.text = obj
    r_txt.font.name = "Times New Roman"
    r_txt.font.size = Pt(16)
    r_txt.font.color.rgb = COLOR_BODY

slide5.shapes.add_picture(r'd:\B.E final year project\icons\chatbot_icon.png', Inches(8.8), Inches(1.9), Inches(0.45), Inches(0.45))


# ----------------------------------------------------
# SLIDE 6: Literature Survey
# ----------------------------------------------------
slide6 = create_base_slide(prs, "Literature Survey")

# 3 Paper Cards in upper region
papers = [
    ("Research Paper 1", "Whisper: Robust Speech Recognition via Large-Scale Weak Supervision", ["• High-accuracy Speech-to-Text", "• Noise-resistant speech recognition"]),
    ("Research Paper 2", "GPT-based Conversational AI for Customer Support", ["• Natural language understanding", "• Context-aware response generation"]),
    ("Research Paper 3", "Text-to-Speech Synthesis Using Neural Networks", ["• Human-like voice generation", "• Low-latency speech synthesis"])
]

card_w = Inches(2.8)
card_h = Inches(1.9)
spacing = Inches(0.2)

for i, (p_num, p_title, p_contribs) in enumerate(papers):
    c_left = Inches(0.6) + i * (card_w + spacing)
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(1.8), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(245, 247, 250)
    card.line.color.rgb = COLOR_NAVY
    card.line.width = Pt(1)
    
    tf = card.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.12)
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = p_num + "\n"
    r.font.name = "Georgia"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = COLOR_NAVY
    
    r_t = p.add_run()
    r_t.text = f'"{p_title}"\n'
    r_t.font.name = "Times New Roman"
    r_t.font.size = Pt(12)
    r_t.font.bold = True
    r_t.font.color.rgb = COLOR_DARK_RED
    
    for cb in p_contribs:
        p_c = tf.add_paragraph()
        p_c.space_before = Pt(2)
        r_c = p_c.add_run()
        r_c.text = cb
        r_c.font.name = "Times New Roman"
        r_c.font.size = Pt(12)
        r_c.font.color.rgb = COLOR_BODY

# Bottom Card: Research Gap
gap_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.85), Inches(8.8), Inches(1.35))
gap_card.fill.solid()
gap_card.fill.fore_color.rgb = RGBColor(254, 243, 235)
gap_card.line.color.rgb = COLOR_RED
gap_card.line.width = Pt(1.5)

tf = gap_card.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.15)
tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run()
r.text = "🔍 Research Gap:"
r.font.name = "Georgia"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = COLOR_RED

gaps = [
    "Existing IVR systems are rigid & rule-based",
    "Limited contextual understanding",
    "Lack of natural conversation flow",
    "High dependency on human agents"
]

p_g = tf.add_paragraph()
p_g.space_before = Pt(4)
r_g = p_g.add_run()
r_g.text = " • " + "   • ".join(gaps)
r_g.font.name = "Times New Roman"
r_g.font.size = Pt(14)
r_g.font.bold = True
r_g.font.color.rgb = COLOR_BODY


# ----------------------------------------------------
# SLIDE 7: Project Overview (Workflow Diagram + Tech Stack + Outcome)
# ----------------------------------------------------
slide7 = create_base_slide(prs, "Project Overview")

# Left Column: Vertical Workflow Diagram
wf_box_w = Inches(3.8)
wf_box_h = Inches(3.4)
wf_container = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), wf_box_w, wf_box_h)
wf_container.fill.solid()
wf_container.fill.fore_color.rgb = RGBColor(245, 247, 252)
wf_container.line.color.rgb = COLOR_NAVY
wf_container.line.width = Pt(1.5)

tf_wf = wf_container.text_frame
tf_wf.margin_left = tf_wf.margin_right = tf_wf.margin_top = tf_wf.margin_bottom = Inches(0.1)

# Header inside workflow container
p_wf_hdr = tf_wf.paragraphs[0]
p_wf_hdr.alignment = PP_ALIGN.CENTER
r = p_wf_hdr.add_run()
r.text = "⚙️ System Workflow"
r.font.name = "Georgia"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = COLOR_NAVY

# Vertical Workflow steps
workflow_steps = [
    "1. User Speaks",
    "2. Speech-to-Text (STT)",
    "3. Speech Processing",
    "4. Intent Detection",
    "5. Large Language Model",
    "6. Response Generation",
    "7. Text-to-Speech (TTS)",
    "8. Voice Response to User"
]

step_top = Inches(2.25)
step_w = Inches(3.4)
step_h = Inches(0.28)

for i, step_text in enumerate(workflow_steps):
    s_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), step_top + i * Inches(0.35), step_w, step_h)
    s_box.fill.solid()
    if i in [0, 7]:
        s_box.fill.fore_color.rgb = COLOR_NAVY
        txt_color = RGBColor(255, 255, 255)
    elif i in [1, 6]:
        s_box.fill.fore_color.rgb = COLOR_MAGENTA
        txt_color = RGBColor(255, 255, 255)
    elif i in [4, 5]:
        s_box.fill.fore_color.rgb = COLOR_RED
        txt_color = RGBColor(255, 255, 255)
    else:
        s_box.fill.fore_color.rgb = RGBColor(230, 235, 245)
        txt_color = COLOR_NAVY
    s_box.line.fill.background()
    
    tf_s = s_box.text_frame
    tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = Inches(0.02)
    p_s = tf_s.paragraphs[0]
    p_s.alignment = PP_ALIGN.CENTER
    r_s = p_s.add_run()
    r_s.text = step_text
    r_s.font.name = "Times New Roman"
    r_s.font.size = Pt(12)
    r_s.font.bold = True
    r_s.font.color.rgb = txt_color

# Right Column Top: Technology Stack
tech_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.8), Inches(4.7), Inches(1.95))
tech_card.fill.solid()
tech_card.fill.fore_color.rgb = COLOR_CARD_BG
tech_card.line.color.rgb = COLOR_CARD_BORDER
tech_card.line.width = Pt(1)

tf_tech = tech_card.text_frame
tf_tech.margin_left = tf_tech.margin_right = tf_tech.margin_top = tf_tech.margin_bottom = Inches(0.12)
tf_tech.word_wrap = True

p_t_hdr = tf_tech.paragraphs[0]
r = p_t_hdr.add_run()
r.text = "💻 Technology Stack"
r.font.name = "Georgia"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = COLOR_NAVY

tech_items = [
    "• Language & Framework: Python, FastAPI",
    "• Speech-to-Text: Whisper API",
    "• Intelligence Engine: GPT / Llama (LLM)",
    "• Text-to-Speech: ElevenLabs Neural TTS",
    "• Database & UI: MongoDB, React (Optional)"
]
for ti in tech_items:
    p = tf_tech.add_paragraph()
    p.space_before = Pt(2)
    r = p.add_run()
    r.text = ti
    r.font.name = "Times New Roman"
    r.font.size = Pt(13)
    r.font.color.rgb = COLOR_BODY

# Right Column Bottom: Expected Outcome
out_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(3.85), Inches(4.7), Inches(1.35))
out_card.fill.solid()
out_card.fill.fore_color.rgb = RGBColor(240, 248, 240) # soft green tint
out_card.line.color.rgb = RGBColor(40, 140, 60)
out_card.line.width = Pt(1.5)

tf_out = out_card.text_frame
tf_out.margin_left = tf_out.margin_right = tf_out.margin_top = tf_out.margin_bottom = Inches(0.12)
tf_out.word_wrap = True

p_o_hdr = tf_out.paragraphs[0]
r = p_o_hdr.add_run()
r.text = "🎯 Expected Outcome"
r.font.name = "Georgia"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = RGBColor(30, 110, 45)

p_o_b = tf_out.add_paragraph()
p_o_b.space_before = Pt(3)
r_o = p_o_b.add_run()
r_o.text = "The proposed system delivers intelligent, real-time, voice-based customer support with faster response time, higher accuracy, reduced operational cost, and improved user experience."
r_o.font.name = "Times New Roman"
r_o.font.size = Pt(13)
r_o.font.color.rgb = COLOR_BODY


# ----------------------------------------------------
# SLIDE 8: Thank You
# ----------------------------------------------------
slide8 = clone_slide(prs, slide_template)

# Clear body placeholder
for sp in list(slide8.shapes):
    if sp.name.startswith("Google Shape;63"): # Title
        tf = sp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ""
    elif sp.name.startswith("Google Shape;64"): # Body
        sp_elem = sp.element
        sp_elem.getparent().remove(sp_elem)

# Center Thank You box
ty_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(7.0), Inches(2.8))
ty_box.fill.solid()
ty_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
ty_box.line.color.rgb = COLOR_NAVY
ty_box.line.width = Pt(2)

tf_ty = ty_box.text_frame
tf_ty.margin_left = tf_ty.margin_right = tf_ty.margin_top = tf_ty.margin_bottom = Inches(0.3)
tf_ty.word_wrap = True

p1 = tf_ty.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run()
r1.text = "Thank You!"
r1.font.name = "Times New Roman"
r1.font.size = Pt(40)
r1.font.bold = True
r1.font.color.rgb = COLOR_RED

p_space = tf_ty.add_paragraph()
p_space.text = ""

p2 = tf_ty.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Questions & Discussion"
r2.font.name = "Georgia"
r2.font.size = Pt(24)
r2.font.bold = True
r2.font.color.rgb = COLOR_NAVY

# Save presentation
out_pptx_path = r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx'
prs.save(out_pptx_path)
print(f"Successfully generated presentation with {len(prs.slides)} slides at {out_pptx_path}!")
