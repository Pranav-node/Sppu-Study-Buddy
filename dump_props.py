import pptx

prs = pptx.Presentation(r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx')

for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    print("Background:", slide.background.fill.type if slide.background else "No background object")
    for shape in slide.shapes:
        print(f"  Shape: name='{shape.name}', type={shape.shape_type}, left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
        if shape.has_text_frame:
            tf = shape.text_frame
            print(f"    TextFrame margins: L={tf.margin_left}, R={tf.margin_right}, T={tf.margin_top}, B={tf.margin_bottom}, word_wrap={tf.word_wrap}")
            for p in tf.paragraphs:
                print(f"    Paragraph level={p.level}, align={p.alignment}: '{p.text}'")
                for r in p.runs:
                    fn = r.font.name
                    fs = r.font.size.pt if r.font.size else None
                    fb = r.font.bold
                    fi = r.font.italic
                    try:
                        fc = r.font.color.rgb if r.font.color.type == 1 else f"type_{r.font.color.type}"
                    except:
                        fc = "none"
                    print(f"      Run: text='{r.text}', font={fn}, sz={fs}, bold={fb}, color={fc}")

