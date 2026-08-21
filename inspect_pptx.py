import pptx
import sys

prs = pptx.Presentation(r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Slide width: {prs.slide_width.inches} inches, height: {prs.slide_height.inches} inches')

for idx, slide in enumerate(prs.slides):
    print(f'\n========================================')
    print(f'SLIDE {idx+1} (Layout: {slide.slide_layout.name})')
    print(f'========================================')
    for shape_idx, shape in enumerate(slide.shapes):
        pos = f"left={shape.left.inches:.2f}\", top={shape.top.inches:.2f}\", width={shape.width.inches:.2f}\", height={shape.height.inches:.2f}\""
        ph_type = shape.placeholder_format.type if shape.is_placeholder else "None"
        print(f'\n  Shape {shape_idx}: "{shape.name}" (Type: {shape.shape_type}, PH: {ph_type}) | {pos}')
        if shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, p in enumerate(tf.paragraphs):
                runs_info = []
                for r in p.runs:
                    font_info = f"font={r.font.name}, size={r.font.size.pt if r.font.size else 'default'}, bold={r.font.bold}, color={r.font.color.rgb if r.font.color and r.font.color.type==1 else r.font.color.type if r.font.color else 'default'}"
                    runs_info.append(f"['{r.text}' ({font_info})]")
                print(f'    P{p_idx} (level={p.level}): "{p.text}" | Runs: {runs_info}')
