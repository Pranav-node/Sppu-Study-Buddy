import pptx

prs = pptx.Presentation(r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx')

print("--- SLIDE MASTER & LAYOUTS ---")
for m_idx, master in enumerate(prs.slide_masters):
    print(f"Master {m_idx}:")
    for s_idx, shape in enumerate(master.shapes):
        print(f"  Master Shape {s_idx}: {shape.name} ({shape.shape_type})")
    for l_idx, layout in enumerate(master.slide_layouts):
        print(f"  Layout {l_idx}: '{layout.name}' with {len(layout.shapes)} shapes")
        for s_idx, shape in enumerate(layout.shapes):
            ph_type = shape.placeholder_format.type if shape.is_placeholder else "None"
            print(f"    Layout Shape {s_idx}: {shape.name} ({shape.shape_type}, PH: {ph_type})")

print("\n--- DETAILED SLIDE SHAPE ANALYSIS ---")
for idx, slide in enumerate(prs.slides):
    print(f"\nSLIDE {idx+1}:")
    for shape_idx, shape in enumerate(slide.shapes):
        print(f" Shape {shape_idx}: {shape.name} | Type: {shape.shape_type}")
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for g_idx, g_shape in enumerate(shape.shapes):
                print(f"   Group Child {g_idx}: {g_shape.name} | Type: {g_shape.shape_type}")
                if g_shape.has_text_frame:
                    for p in g_shape.text_frame.paragraphs:
                        print(f"     Text: '{p.text}'")
        elif shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, r in enumerate(p.runs):
                    r_color = r.font.color.rgb if r.font.color and r.font.color.type == 1 else "not-rgb"
                    print(f"    P{p_idx} R{r_idx}: text='{r.text}' | font={r.font.name}, sz={r.font.size.pt if r.font.size else None}, bold={r.font.bold}, italic={r.font.italic}, color={r_color}")

