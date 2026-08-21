import pptx

prs = pptx.Presentation(r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx')

print("--- SLIDE 1 SHAPES ---")
for s in prs.slides[0].shapes:
    print(s.name, s.shape_type, s.left, s.top, s.width, s.height)

print("--- SLIDE 2 SHAPES ---")
for s in prs.slides[1].shapes:
    print(s.name, s.shape_type, s.left, s.top, s.width, s.height)

print("--- SLIDE 3 SHAPES ---")
for s in prs.slides[2].shapes:
    print(s.name, s.shape_type, s.left, s.top, s.width, s.height)
