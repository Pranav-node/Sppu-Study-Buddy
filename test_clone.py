import pptx
from pptx.oxml import parse_xml

def clone_slide(prs, source_slide):
    # Create a new slide using the same layout as source_slide
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Remove default shapes on new slide if any
    for sp in list(new_slide.shapes):
        sp_elem = sp.element
        sp_elem.getparent().remove(sp_elem)
        
    # Copy all shapes from source_slide to new_slide
    for sp in source_slide.shapes:
        new_sp = parse_xml(sp.element.xml)
        new_slide.shapes._spTree.append(new_sp)
        
    return new_slide

prs = pptx.Presentation(r'd:\B.E final year project\B.E. Project Topic Finalisation Review Template.pptx')
slide2 = prs.slides[1]
new_s = clone_slide(prs, slide2)
prs.save(r'd:\B.E final year project\test_clone.pptx')
print("Successfully cloned slide! Total slides now:", len(prs.slides))
